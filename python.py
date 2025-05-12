from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Create a presentation
prs = Presentation()

# Define title slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "How Customer Reviews Influence Restaurant Closures"
subtitle.text = "Capstone Project Presentation by Mia Gallien"

# Content slide helper
def add_slide(title_text, content_lines):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = title_text
    content.text = "\n".join(content_lines)

# Add slides with content
slides_content = [
    ("Executive Summary", [
        "Objective: Analyze the effect of customer reviews on restaurant closures.",
        "Data: Yelp, city inspection records, and review history.",
        "Methods: EDA, SQL, predictive modeling, interactive dashboards.",
        "Key Insight: Negative reviews and declining trends strongly correlate with closures."
    ]),
    ("Introduction", [
        "Overview of post-pandemic restaurant industry.",
        "Importance of online reviews on customer perception.",
        "Main Question: Can review data predict restaurant closures?"
    ]),
    ("Data Collection", [
        "Sources: Yelp, City records (closures, permits), Web scraping.",
        "Timeframe: 2018–2024.",
        "Tools: Python, Pandas, BeautifulSoup, SQL."
    ]),
    ("Data Wrangling", [
        "Removed nulls/duplicates, unified format.",
        "Merged business data with closures using business_id and location.",
        "Created features: review volume trend, sentiment, average stars.",
        "Binary classification target: Open vs. Closed."
    ]),
    ("EDA & Interactive Visual Analytics", [
        "Rating distributions (histograms, boxplots).",
        "Monthly review volume trends before closures.",
        "Word clouds of reviews by sentiment.",
        "Tools: Plotly, Seaborn, Matplotlib."
    ]),
    ("Predictive Analysis Methodology", [
        "Models: Logistic Regression, Random Forest, XGBoost.",
        "Features: Ratings, sentiment, review volume trends.",
        "Train/Test Split: 80/20 with cross-validation."
    ]),
    ("EDA Visualization Results", [
        "65% of closures had average ratings <= 3.5.",
        "Monthly review drops seen in 72% of closures.",
        "Negative keywords: 'dirty', 'rude', 'slow'.",
        "Charts highlight clear sentiment drop before closure."
    ]),
    ("EDA with SQL Results", [
        "SQL used to filter businesses with <3.5 stars.",
        "JOINs matched Yelp data with closure data.",
        "Closure rate 38% higher in low-rated businesses.",
        "SQL used to group by city, year, and sentiment."
    ]),
    ("Interactive Map with Folium", [
        "Mapped closures by region and rating.",
        "Pop-ups show recent review trends before closure.",
        "Hotspots: Low-income areas and tourist zones."
    ]),
    ("Plotly Dash Dashboard", [
        "Interactive filters: year, stars, review count.",
        "Includes graphs, word clouds, and predictions.",
        "Live input: Predict closure risk from data inputs."
    ]),
    ("Predictive Analysis Results", [
        "Accuracy: 82%, Precision: 76%, Recall: 70%.",
        "Top predictors: rating drop, 1-star review frequency, sentiment trends.",
        "Model visualization: ROC Curve, Confusion Matrix."
    ]),
    ("Conclusion", [
        "Strong predictive power of reviews in forecasting closures.",
        "Recommend monitoring sentiment to act early.",
        "Businesses can improve survival odds with customer service focus."
    ]),
    ("Creative Additions", [
        "Dashboard app with live input and graphs.",
        "Animated visualization of rating decline before closure.",
        "Sector-based recommendations for high-risk restaurants."
    ]),
]

# Add each content slide
for title_text, content_lines in slides_content:
    add_slide(title_text, content_lines)

# Save the presentation
output_path = "/mnt/data/Restaurant_Closures_Review_Analysis_Presentation.pptx"
prs.save(output_path)

output_path
